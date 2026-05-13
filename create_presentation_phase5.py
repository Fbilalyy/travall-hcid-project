"""
IE48L Phase 5A & 6A — Prototype & Field Trial-1 Presentation (PPTX)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
import os

OUT_DIR = "/Users/fatihbilalyilmaz/PycharmProjects/IE492/IE48L"

DARK_BLUE  = RGBColor(0x1E, 0x3A, 0x5F)
MED_BLUE   = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
TEAL       = RGBColor(0x14, 0xB8, 0xA6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
G700       = RGBColor(0x37, 0x41, 0x51)
G500       = RGBColor(0x6B, 0x72, 0x80)
G400       = RGBColor(0x9C, 0xA3, 0xAF)
G100       = RGBColor(0xF3, 0xF4, 0xF6)
RED        = RGBColor(0xDC, 0x26, 0x26)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
ORANGE     = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_RED  = RGBColor(0xFE, 0xF2, 0xF2)
LIGHT_GREEN= RGBColor(0xEC, 0xFC, 0xF2)
LIGHT_ORANGE = RGBColor(0xFF, 0xFB, 0xEB)
LIGHT_TEAL = RGBColor(0xF0, 0xFD, 0xFA)
ACCENT     = RGBColor(0xE5, 0x3E, 0x3E)


def no_border(table):
    tbl = table._tbl
    for cell in tbl.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}tc'):
        tcPr = cell.find('{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr')
        if tcPr is None:
            tcPr = etree.SubElement(cell, '{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr')
        for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
            ln = tcPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}' + border_name)
            if ln is None:
                ln = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}' + border_name)
            ln.set('w', '0')
            nf = ln.find('{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')
            if nf is None:
                etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')


def set_cell_border(table, row, col, sides=('lnB',), color=RGBColor(0xE2, 0xE8, 0xF0), width=12700):
    tbl = table._tbl
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    rows = list(tbl.iter(ns + 'tr'))
    cells = list(rows[row].iter(ns + 'tc'))
    cell = cells[col]
    tcPr = cell.find(ns + 'tcPr')
    if tcPr is None:
        tcPr = etree.SubElement(cell, ns + 'tcPr')
    for side in sides:
        ln = tcPr.find(ns + side)
        if ln is None:
            ln = etree.SubElement(tcPr, ns + side)
        ln.set('w', str(width))
        for child in list(ln):
            ln.remove(child)
        solidFill = etree.SubElement(ln, ns + 'solidFill')
        srgbClr = etree.SubElement(solidFill, ns + 'srgbClr')
        srgbClr.set('val', '%02X%02X%02X' % (color[0] if isinstance(color, tuple) else int(str(color)[:2], 16),
                                                color[1] if isinstance(color, tuple) else int(str(color)[2:4], 16),
                                                color[2] if isinstance(color, tuple) else int(str(color)[4:6], 16)))


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_cell(cell, text, sz=14, bold=False, color=G700, align=PP_ALIGN.LEFT, bg_color=None):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(2)
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color


def set_cell_multi(cell, lines_data, bg_color=None):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, item in enumerate(lines_data):
        text = item[0]
        sz = item[1] if len(item) > 1 else 14
        bold = item[2] if len(item) > 2 else False
        color = item[3] if len(item) > 3 else G700
        align = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(3)
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color


def add_title(slide, text, left=0.5, top=0.2, width=12.3, height=0.7, sz=34):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    return tb


def add_subtitle_badge(slide, text, left=0.5, top=0.85, width=4, height=0.45):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    # We'll just use colored text since shape fill on textbox is tricky
    p.font.color.rgb = ACCENT
    return tb


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: TITLE =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, DARK_BLUE)

    tb = sl.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.3), Inches(6.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "IE48L"
    p.font.size = Pt(24)
    p.font.color.rgb = G400
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(15)

    p = tf.add_paragraph()
    p.text = "Phase 5A & 6A"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "Prototype & Field Trial-1"
    p.font.size = Pt(30)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Project: TravAll"
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(50)

    p = tf.add_paragraph()
    p.text = "Ozde Avdan  |  Nilsu Calliogullari  |  Deniz Tas  |  Fatih Bilal Yilmaz"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "13.04.2026"
    p.font.size = Pt(16)
    p.font.color.rgb = G500
    p.alignment = PP_ALIGN.CENTER

    # ===== SLIDE 2: AGENDA =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "Agenda")

    table_shape = sl.shapes.add_table(7, 2, Inches(1.5), Inches(1.3), Inches(10.3), Inches(5.5))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(0.8)
    table.columns[1].width = Inches(9.5)

    agenda_items = [
        ("1", "Creation of the Prototype"),
        ("2", "Field Trial-1 Tasks"),
        ("3", "Field Trial-1 Results"),
        ("4", "Participant Analysis"),
        ("5", "Feedbacks from Field Trial-1"),
        ("6", "Feedbacks from Phase 4 Presentation"),
        ("7", "Design Modifications"),
    ]
    for i, (num, text) in enumerate(agenda_items):
        set_cell(table.cell(i, 0), num, 20, True, ACCENT, PP_ALIGN.CENTER, G100)
        set_cell(table.cell(i, 1), text, 20, False, G700, PP_ALIGN.LEFT, G100)

    # ===== SLIDE 3: PROTOTYPE CREATION - FIGMA =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "1. Prototype Creation")
    add_subtitle_badge(sl, "Design Phase: Figma")

    table_shape = sl.shapes.add_table(4, 1, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(12.3)

    items = [
        "All screens designed in Figma following lo-fi prototyping principles",
        "Solid colors, emoji-based visuals, repeating component structures",
        "17 screens total: Passport Hub, Country Details, Visa Application, Travel Planner, Exploration Flow, Profile",
        "Screens exported as PNG files and uploaded to Claude with a navigation map (button targets, back navigation, active tabs)",
    ]
    for i, item in enumerate(items):
        set_cell_multi(table.cell(i, 0), [
            ("  " + item, 18, False, G700),
        ], G100 if i % 2 == 0 else WHITE)

    # ===== SLIDE 4: CLICKABLE PROTOTYPE =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "1. Prototype Creation")
    add_subtitle_badge(sl, "Clickable Prototype via Claude")

    table_shape = sl.shapes.add_table(2, 2, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(6.15)
    table.columns[1].width = Inches(6.15)

    set_cell_multi(table.cell(0, 0), [
        ("What Claude Built", 18, True, DARK_BLUE),
        ("", 6),
        ("  17 screens within an actual phone frame", 15, False, G700),
        ("  JavaScript navigation system", 15, False, G700),
        ("  Browsing history logic (back button)", 15, False, G700),
        ("  Bottom navigation bar (4 main tabs)", 15, False, G700),
    ], G100)

    set_cell_multi(table.cell(0, 1), [
        ("Advantages over Flipbook", 18, True, GREEN),
        ("", 6),
        ("  Participants interact as with a real app", 15, False, G700),
        ("  Instant page transitions", 15, False, G700),
        ("  Smooth navigation flow", 15, False, G700),
        ("  Usability data closer to real-life behavior", 15, False, G700),
    ], LIGHT_GREEN)

    set_cell_multi(table.cell(1, 0), [
        ("A digitally transferred and greatly improved variant of lo-fi prototyping in the context of usability testing.", 16, True, DARK_BLUE),
    ], WHITE)
    set_cell_multi(table.cell(1, 1), [("", 8)], WHITE)

    # ===== SLIDE 5: TASKS 1-8 =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "2. Field Trial-1 Tasks")

    table_shape = sl.shapes.add_table(8, 2, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.8))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(12.1)

    tasks_1 = [
        ("1", "Install TravAll, verify passport info, find visa-free country count"),
        ("2", "Plan a trip to Japan — check visa requirements & entry info"),
        ("3", "Go to United Kingdom — determine how to apply for a visa"),
        ("4", "Apply for a UK Standard Visitor visa (destination, type, date)"),
        ("5", "Locate where to renew your expiring Schengen visa"),
        ("6", "Find and turn on the visa expiry reminder notification"),
        ("7", "Plan a Barcelona trip: departure, return, transport, accommodation"),
        ("8", "View your ongoing Tokyo trip schedule (flight, hotel, activities)"),
    ]
    for i, (num, text) in enumerate(tasks_1):
        bg = G100 if i % 2 == 0 else WHITE
        set_cell(table.cell(i, 0), num, 14, True, ACCENT, PP_ALIGN.CENTER, bg)
        set_cell(table.cell(i, 1), text, 14, False, G700, PP_ALIGN.LEFT, bg)

    # ===== SLIDE 6: TASKS 9-15 =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "2. Field Trial-1 Tasks")

    table_shape = sl.shapes.add_table(7, 2, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.8))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(12.1)

    tasks_2 = [
        ("9", "Locate a payment pending item (TeamLab Planets) in Tokyo schedule"),
        ("10", "Click on Explore and read the community feed"),
        ("11", "Open Alex K.'s Tokyo budget breakdown, find accommodation cost"),
        ("12", "Start planning a new trip from the Tokyo route post"),
        ("13", "Visit Profile — see continents explored and proportions"),
        ("14", "Switch off community interaction notifications (likes & comments)"),
        ("15", "Navigate to Passport & Visas settings, find 'Add Second Passport'"),
    ]
    for i, (num, text) in enumerate(tasks_2):
        bg = G100 if i % 2 == 0 else WHITE
        set_cell(table.cell(i, 0), num, 14, True, ACCENT, PP_ALIGN.CENTER, bg)
        set_cell(table.cell(i, 1), text, 14, False, G700, PP_ALIGN.LEFT, bg)

    # ===== SLIDE 7: RESULTS TABLE =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "3. Field Trial-1 Results")
    add_subtitle_badge(sl, "Per-Participant Results")

    table_shape = sl.shapes.add_table(9, 6, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(3.8)
    for c in range(1, 6):
        table.columns[c].width = Inches(1.7)

    # Header
    headers = ["Metric", "P1", "P2", "P3", "P4", "P5"]
    for c, h in enumerate(headers):
        set_cell(table.cell(0, c), h, 14, True, WHITE, PP_ALIGN.CENTER, DARK_BLUE)

    data = [
        ("Completion time",        ["6' 12\"", "7' 45\"", "5' 58\"", "8' 03\"", "6' 34\""]),
        ("Task completion rate",   ["100%", "93%", "100%", "87%", "100%"]),
        ("Number of errors",       ["2", "3", "1", "4", "1"]),
        ("Back nav needed",        ["3", "5", "2", "6", "3"]),
        ("Likability",             ["9", "8", "9", "7", "10"]),
        ("Ease of use",            ["8", "7", "9", "6", "9"]),
        ("Visual clarity",         ["9", "9", "10", "8", "9"]),
        ("Navigation comfort",    ["8", "7", "9", "7", "10"]),
    ]
    for r, (metric, vals) in enumerate(data):
        bg = G100 if r % 2 == 0 else WHITE
        set_cell(table.cell(r+1, 0), metric, 13, True, G700, PP_ALIGN.LEFT, bg)
        for c, v in enumerate(vals):
            color = G700
            if metric == "Task completion rate" and v != "100%":
                color = RED
            elif metric == "Number of errors" and int(v) >= 4:
                color = RED
            set_cell(table.cell(r+1, c+1), v, 13, False, color, PP_ALIGN.CENTER, bg)

    # ===== SLIDE 8: AGGREGATE METRICS =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "3. Field Trial-1 Results")
    add_subtitle_badge(sl, "Aggregate Metrics")

    table_shape = sl.shapes.add_table(7, 2, Inches(2.0), Inches(1.5), Inches(9.3), Inches(5.0))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(6.5)
    table.columns[1].width = Inches(2.8)

    agg = [
        ("Average task completion time", "6' 54\""),
        ("Overall task completion success rate", "96%"),
        ("Average number of errors per user", "2.2"),
        ("Average likability score", "8.6 / 10"),
        ("Average ease of use score", "7.8 / 10"),
        ("Average visual clarity score", "9.0 / 10"),
        ("Average navigation comfort score", "8.2 / 10"),
    ]
    for i, (metric, val) in enumerate(agg):
        bg = G100 if i % 2 == 0 else WHITE
        set_cell(table.cell(i, 0), metric, 16, False, G700, PP_ALIGN.LEFT, bg)
        set_cell(table.cell(i, 1), val, 16, True, DARK_BLUE, PP_ALIGN.CENTER, bg)

    # ===== SLIDE 9: KEY OBSERVATIONS =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "3. Field Trial-1 Results")
    add_subtitle_badge(sl, "Key Observations")

    table_shape = sl.shapes.add_table(6, 1, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(12.3)

    observations = [
        "Average completion time: 6-8 minutes across all participants",
        "Longer times on visa application and country detail screens (participants reading content, not usability issues)",
        "Participant 4 made the most mistakes, mostly on 'Renew Visa' and 'Notification Settings' flows",
        "Task 4 (visa form) and Task 10 (Explore feed navigation) were most time-consuming",
        "Back navigation button was intuitive — participants quickly recovered with it",
        "Visual clarity rated highest (9.0/10), ease of use rated lowest (7.8/10)",
    ]
    for i, obs in enumerate(observations):
        bg = G100 if i % 2 == 0 else WHITE
        set_cell(table.cell(i, 0), "  " + obs, 16, False, G700, PP_ALIGN.LEFT, bg)

    # ===== SLIDE 10: PARTICIPANT 2 =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "4. Participant Analysis")

    # P2 header
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Participant 2 — 93% Success Rate (1 Incomplete Task)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    table_shape = sl.shapes.add_table(2, 1, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(12.3)

    set_cell_multi(table.cell(0, 0), [
        ("Failed Task 6: Could not find notification settings", 16, True, RED),
        ("", 6),
        ("  Tapped through bottom nav bar (Visa, My Trips, Explore) until reaching Profile", 14, False, G700),
        ("  Wandered past stats section and world map before finding Account subsection", 14, False, G700),
        ("  Tapped continent progress bars twice, hoping they were interactive", 14, False, G700),
        ("  At ~90 seconds, orally stated she could not locate settings", 14, False, G700),
    ], LIGHT_ORANGE)

    set_cell_multi(table.cell(1, 0), [
        ("Post-Trial Interview:", 14, True, DARK_BLUE),
        ("", 4),
        ("\"I assumed there must be some sort of gear icon or a tab of settings somewhere", 14, False, G500),
        ("maybe in the top right of the profile. I did not know that the list at the bottom", 14, False, G500),
        ("of the profile page was the settings menu.\"", 14, False, G500),
        ("", 6),
        ("Insight: Settings affordance not clearly differentiated from informational content.", 14, True, DARK_BLUE),
        ("Missing gear icon breaks mental model (Instagram, Spotify convention).", 14, False, G700),
    ], G100)

    # ===== SLIDE 11: PARTICIPANT 4 =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "4. Participant Analysis")

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Participant 4 — 87% Success Rate (2 Incomplete Tasks)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED

    table_shape = sl.shapes.add_table(2, 1, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(12.3)

    set_cell_multi(table.cell(0, 0), [
        ("Failed Task 4: Tapped 'Renew Visa' instead of 'Apply for Visa'", 15, True, RED),
        ("", 4),
        ("  Sought a form on Renew Visa screen, couldn't find one. Tapped back, tried Japan (visa on arrival).", 13, False, G700),
        ("  \"The two bottom buttons seemed similar to me, I did not read them properly at first", 13, False, G500),
        ("  and simply pressed the second one, perhaps they should have been more different.\"", 13, False, G500),
    ], LIGHT_RED)

    set_cell_multi(table.cell(1, 0), [
        ("Failed Task 15: Found 'Add Second Passport' card but ignored it", 15, True, RED),
        ("", 4),
        ("  Dashed border appeared decorative, light blue color made it blend in.", 13, False, G700),
        ("  Tapped Turkish Passport card instead, searched for an add option for 2+ minutes.", 13, False, G700),
        ("  \"That was no box to tap — that was more of a place-bosom than a box.\"", 13, False, G500),
        ("", 4),
        ("  Insight: Dashed border + light color = perceived as placeholder/decorative, not interactive.", 13, True, DARK_BLUE),
    ], LIGHT_ORANGE)

    # ===== SLIDE 12: POSITIVE FEEDBACK =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "5. Feedbacks from Field Trial-1")
    add_subtitle_badge(sl, "Positive Feedback")

    table_shape = sl.shapes.add_table(6, 1, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(12.3)

    positives = [
        "Blue color scheme feels travel-oriented and trustworthy, familiar to airline apps",
        "Visa status badges (green, amber, red) on home screen make status instantly visible",
        "Tokyo trip timeline (flight, hotel, activity on single page) is exactly what users want in a travel app",
        "Passport detail page with global ranking and visa-free country breakdown was a pleasant surprise",
        "Four-tab navigation kept users oriented — never felt lost in the app",
        "Community feed was a pleasant surprise, would use it before making travel plans",
    ]
    for i, fb in enumerate(positives):
        set_cell(table.cell(i, 0), "  + " + fb, 15, False, GREEN, PP_ALIGN.LEFT, LIGHT_GREEN if i % 2 == 0 else WHITE)

    # ===== SLIDE 13: NEGATIVE FEEDBACK =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "5. Feedbacks from Field Trial-1")
    add_subtitle_badge(sl, "Negative Feedback")

    table_shape = sl.shapes.add_table(6, 1, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(12.3)

    negatives = [
        "'Apply for Visa' vs 'Renew Visa' buttons look too similar, caused wrong selection",
        "Settings section blends into Profile content — discovered late as a menu",
        "'Add Second Passport' dashed border resembled a placeholder, not a tappable element",
        "Passport Hub has too much information — passport card, search box, and country list compete",
        "Community feed lacks filtering by destination or budget range",
        "World map on Profile page is a dummy box — feels incomplete vs. other sections",
    ]
    for i, fb in enumerate(negatives):
        set_cell(table.cell(i, 0), "  - " + fb, 15, False, RED, PP_ALIGN.LEFT, LIGHT_RED if i % 2 == 0 else WHITE)

    # ===== SLIDE 14: PHASE 4 FEEDBACK =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "6. Feedbacks from Phase 4 Presentation")
    add_subtitle_badge(sl, "The Fidelity Question")

    table_shape = sl.shapes.add_table(2, 1, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(12.3)

    set_cell_multi(table.cell(0, 0), [
        ("Panel Concerns on Initial Designs", 17, True, DARK_BLUE),
        ("", 6),
        ("  Custom iconography, layered card forms, mixed photographic placeholders", 15, False, G700),
        ("  Elaborate micro-interactions — too detailed and visually dense for a lo-fi phase", 15, False, RED),
        ("  The amount of detail might confuse lo-fi and hi-fi deliverables", 15, False, G700),
        ("  Harder to evaluate underlying UX decisions separately", 15, False, G700),
    ], LIGHT_RED)

    set_cell_multi(table.cell(1, 0), [
        ("What the Panel Recommended", 17, True, DARK_BLUE),
        ("", 6),
        ("  Main goal at this stage: convey structural rationale, user flow, and information hierarchy", 15, False, G700),
        ("  Focus on refinement of interaction logic, not visual polish", 15, False, G700),
        ("  The team engaged in careful and systematic modifications before Field Trial-1", 15, True, GREEN),
    ], LIGHT_GREEN)

    # ===== SLIDE 15: DESIGN MOD 1 =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "7. Design Modifications")

    table_shape = sl.shapes.add_table(3, 2, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(6.15)
    table.columns[1].width = Inches(6.15)

    # Row 0: Headers
    set_cell(table.cell(0, 0), "Abandon Elaborate Visuals", 18, True, RED, PP_ALIGN.CENTER, LIGHT_RED)
    set_cell(table.cell(0, 1), "Simplification of Screen Logic", 18, True, MED_BLUE, PP_ALIGN.CENTER, LIGHT_BLUE)

    # Row 1: Before
    set_cell_multi(table.cell(1, 0), [
        ("BEFORE", 13, True, RED),
        ("  Mixed placeholder images", 13, False, G700),
        ("  Abstract icon sets", 13, False, G700),
        ("  Custom-drawn UI elements", 13, False, G700),
        ("  Created visual noise", 13, False, G500),
        ("", 6),
        ("AFTER", 13, True, GREEN),
        ("  Emoji-like, globally familiar symbols", 13, False, G700),
        ("  No visual decoding needed", 13, False, G700),
        ("  Universally comprehensible", 13, False, G700),
    ], G100)

    set_cell_multi(table.cell(1, 1), [
        ("BEFORE", 13, True, RED),
        ("  Multiple embedded interaction layers", 13, False, G700),
        ("  Gradient backgrounds, shadows", 13, False, G700),
        ("  Decorative card treatments", 13, False, G700),
        ("", 6),
        ("AFTER", 13, True, GREEN),
        ("  One major action per screen", 13, False, G700),
        ("  Flat color fills, uniform blue ramp", 13, False, G700),
        ("  Low border weights, wide whitespace", 13, False, G700),
        ("  Easier to trace errors to specific choices", 13, False, G700),
    ], G100)

    # Row 2: Component standardization
    set_cell_multi(table.cell(2, 0), [
        ("Component Standardization", 16, True, DARK_BLUE),
        ("  Consistent reusable elements:", 13, False, G700),
        ("  Top nav bar, badge system, card structure,", 13, False, G700),
        ("  timeline element, bottom tab bar", 13, False, G700),
    ], WHITE)

    set_cell_multi(table.cell(2, 1), [
        ("Result", 16, True, GREEN),
        ("  Less time reorienting between screens", 13, False, G700),
        ("  More dependable usability data", 13, False, G700),
        ("  Visual consistency isolated from behavior", 13, False, G700),
    ], LIGHT_GREEN)

    # ===== SLIDE 16: OUTCOME =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)
    add_title(sl, "7. Design Modifications — Outcome")

    table_shape = sl.shapes.add_table(1, 3, Inches(1.0), Inches(1.5), Inches(11.3), Inches(2.0))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(3.77)
    table.columns[1].width = Inches(3.77)
    table.columns[2].width = Inches(3.77)

    set_cell_multi(table.cell(0, 0), [
        ("Structurally Complete", 16, True, DARK_BLUE),
        ("", 4),
        ("All 17 screens with full", 14, False, G700),
        ("navigation implemented", 14, False, G700),
    ], G100)

    set_cell_multi(table.cell(0, 1), [
        ("Visually Decipherable", 16, True, DARK_BLUE),
        ("", 4),
        ("Emoji-based, clean,", 14, False, G700),
        ("no visual noise", 14, False, G700),
    ], G100)

    set_cell_multi(table.cell(0, 2), [
        ("Neutral Enough", 16, True, DARK_BLUE),
        ("", 4),
        ("Participants respond to", 14, False, G700),
        ("IA, not aesthetics", 14, False, G700),
    ], G100)

    # Evidence
    table_shape2 = sl.shapes.add_table(3, 2, Inches(2.5), Inches(4.0), Inches(8.3), Inches(2.5))
    table2 = table_shape2.table
    no_border(table2)
    table2.columns[0].width = Inches(5.5)
    table2.columns[1].width = Inches(2.8)

    evidence = [
        ("Visual clarity score", "9.0 / 10"),
        ("Task completion success rate", "96%"),
        ("Navigation comfort score", "8.2 / 10"),
    ]
    for i, (label, val) in enumerate(evidence):
        bg = LIGHT_GREEN if i % 2 == 0 else WHITE
        set_cell(table2.cell(i, 0), label, 16, False, G700, PP_ALIGN.LEFT, bg)
        set_cell(table2.cell(i, 1), val, 18, True, GREEN, PP_ALIGN.CENTER, bg)

    # ===== SLIDE 17: THANK YOU =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, DARK_BLUE)

    tb = sl.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(30)

    p = tf.add_paragraph()
    p.text = "Questions & Discussion"
    p.font.size = Pt(26)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(40)

    p = tf.add_paragraph()
    p.text = "Ozde  |  Nilsu  |  Deniz  |  Fatih Bilal"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    p.alignment = PP_ALIGN.CENTER

    out = os.path.join(OUT_DIR, "IE48L_Phase5A_Presentation.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    return out


if __name__ == '__main__':
    print("Building Phase 5A & 6A presentation...")
    build()
    print("Done!")
