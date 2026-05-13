"""
5-Minute Presentation — IE48L Phase 3A & 3B (8 slides)
Uses TABLES instead of floating shapes — no layer shifting
"""
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.table import _Cell
import os

PDF_PATH = "/Users/fatihbilalyilmaz/PycharmProjects/IE492/IE48L/IE48L-Project Phase 3 (1).pdf"
OUT_DIR = "/Users/fatihbilalyilmaz/PycharmProjects/IE492/IE48L"

DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
MED_BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
G700 = RGBColor(0x37, 0x41, 0x51)
G500 = RGBColor(0x6B, 0x72, 0x80)
G100 = RGBColor(0xF3, 0xF4, 0xF6)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_RED = RGBColor(0xFE, 0xF2, 0xF2)
LIGHT_GREEN = RGBColor(0xEC, 0xFC, 0xF2)


def extract_images(pdf_path, out_dir):
    doc = fitz.open(pdf_path)
    imgs = []
    for idx in [8, 9, 10]:
        page = doc[idx]
        pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))
        p = os.path.join(out_dir, f"sb_{idx-7}.png")
        pix.save(p)
        imgs.append(p)
    doc.close()
    return imgs


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_cell(cell, text, sz=14, bold=False, color=G700, align=PP_ALIGN.LEFT, bg_color=None):
    """Set cell text and style"""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(2)

    if bg_color:
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = bg_color


def set_cell_multi(cell, lines_data, bg_color=None):
    """Set cell with multiple styled lines: [(text, sz, bold, color), ...]"""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)

    for i, item in enumerate(lines_data):
        text = item[0]
        sz = item[1] if len(item) > 1 else 14
        bold = item[2] if len(item) > 2 else False
        color = item[3] if len(item) > 3 else G700
        align = item[4] if len(item) > 4 else PP_ALIGN.LEFT

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(3)

    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color


def no_border(table):
    """Remove all table borders"""
    from lxml import etree
    tbl = table._tbl
    for cell in tbl.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}tc'):
        tcPr = cell.find('{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr')
        if tcPr is None:
            tcPr = etree.SubElement(cell, '{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr')
        for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
            ln = tcPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}' + border_name)
            if ln is None:
                ln = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}' + border_name)
            ln.set('w', '0')
            noFill = ln.find('{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')
            if noFill is None:
                etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}noFill')


def build(imgs):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: TITLE =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, DARK_BLUE)

    # Single textbox for title content - no layers
    tb = sl.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11.3), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "IE48L — Phase 3A & 3B"
    p.font.size = Pt(46); p.font.bold = True; p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(20)

    p = tf.add_paragraph()
    p.text = "Persona Creation & Sample Narrative Storyboard"
    p.font.size = Pt(28); p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(50)

    p = tf.add_paragraph()
    p.text = "Özde Avdan  •  Nilsu Çallıoğulları  •  Deniz Taş  •  Fatih Bilal Yılmaz"
    p.font.size = Pt(18); p.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "18.03.2026"
    p.font.size = Pt(16); p.font.color.rgb = G500
    p.alignment = PP_ALIGN.CENTER

    # ===== SLIDE 2: PROBLEM & TRAVALL =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)

    # Title textbox
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Problem & TravAll's Vision"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(12)
    p = tf.add_paragraph()
    p.text = "Modern travelers juggle 3-4 apps, face visa anxiety, and waste time on manual data entry. TravAll is a super-app that consolidates everything into one seamless experience."
    p.font.size = Pt(18); p.font.color.rgb = G700

    # Stats as TABLE (1 row, 4 cols) - stable!
    table_shape = sl.shapes.add_table(2, 4, Inches(0.5), Inches(2.8), Inches(12.3), Inches(3.0))
    table = table_shape.table
    no_border(table)

    stats = [
        ("58.49%", "face visa uncertainty", RED),
        ("69.81%", "suffer from document clutter", ORANGE),
        ("81.13%", "want a simple barebones UI", MED_BLUE),
        ("94.34%", "need final approval control", GREEN),
    ]
    for i, (val, label, c) in enumerate(stats):
        table.columns[i].width = Inches(3.075)
        # Value row
        set_cell(table.cell(0, i), val, 38, True, c, PP_ALIGN.CENTER, G100)
        # Label row
        set_cell(table.cell(1, i), label, 16, False, G700, PP_ALIGN.CENTER, G100)

    # Source
    tb2 = sl.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4))
    p = tb2.text_frame.paragraphs[0]
    p.text = "Source: Phase 2 Survey Results"
    p.font.size = Pt(12); p.font.color.rgb = G500; p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    # ===== SLIDE 3: THREE PERSONAS =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Three Personas"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = DARK_BLUE

    # Personas as TABLE (6 rows, 3 cols)
    table_shape = sl.shapes.add_table(6, 3, Inches(0.3), Inches(1.0), Inches(12.7), Inches(6.2))
    table = table_shape.table
    no_border(table)

    for i in range(3):
        table.columns[i].width = Inches(4.23)

    personas = [
        ("🎒  Yavuz Polat", MED_BLUE),
        ("💼  Nihan Katırcı", TEAL),
        ("😰  Melih Özgen", ORANGE),
    ]
    titles = [
        "The Efficient Traveller",
        "The High Frequency Business Voyager",
        "The Anxious First Timer",
    ]
    demos = [
        "24, Male, Chemistry Senior at Boğaziçi\nBurgundy Passport • Pfizer intern",
        "38, Female, MBA, C-Level Executive\nGreen Passport • 12+ trips/year",
        "21, Male, Undergraduate Student\nOrdinary (Bordo) Passport • Low budget",
    ]
    features = [
        "• One-handed smartphone navigation\n• Minimalist interface preference\n• D&D Dungeon Master, cocktail culture",
        "• Wants a 'Super App' (replace 3-4 apps)\n• Human-like AI communication\n• Luxury travel & fine dining",
        "• Visual learner (icons over text)\n• First independent trip abroad\n• Photography & social media",
    ]
    pains = [
        "Pain Points:\n• App fragmentation kills motivation\n• Manual calendar = Muda (waste)\n• Documents scattered across systems",
        "Pain Points:\n• Border rules change without notice\n• Can't find tickets in 3-4 apps\n• Price transparency is a deal breaker",
        "Pain Points:\n• Visa research = paralysis & anxiety\n• Might cancel trip (uncertainty)\n• Demands final decision control",
    ]

    for i in range(3):
        c = personas[i][1]
        # Row 0: Name
        set_cell(table.cell(0, i), personas[i][0], 20, True, c, PP_ALIGN.CENTER, G100)
        # Row 1: Title
        set_cell(table.cell(1, i), titles[i], 12, False, G500, PP_ALIGN.CENTER, G100)
        # Row 2: Demographics
        set_cell(table.cell(2, i), demos[i], 12, False, G700, PP_ALIGN.LEFT, WHITE)
        # Row 3: Features
        set_cell(table.cell(3, i), features[i], 11, False, G700, PP_ALIGN.LEFT, WHITE)
        # Row 4: Pain points
        set_cell(table.cell(4, i), pains[i], 11, False, G700, PP_ALIGN.LEFT, LIGHT_RED)
        # Row 5: spacer
        set_cell(table.cell(5, i), "", 4, False, G700, PP_ALIGN.LEFT, WHITE)

    # ===== SLIDES 4-6: SCENARIO + STORYBOARD =====
    scenarios = [
        ("Scenario 1: Yavuz Polat", MED_BLUE,
         "❌  BEFORE TravAll\n\n• Ticket in Gmail, hotel in Booking\n• Visa research on consulate sites\n• Manual calendar entry (Muda!)\n• 3+ apps juggling to the airport",
         "✅  AFTER TravAll\n\n• All tickets auto-imported to wallet\n• Visual pictorial visa flow\n• One-handed navigation on subway\n• Single QR at passport control",
         imgs[0]),
        ("Scenario 2: Nihan Katırcı", TEAL,
         "❌  BEFORE TravAll\n\n• Border rules shift without warning\n• Documents in 3-4 different apps\n• Can't find ticket before meeting\n• Stuck in queue — rule changed",
         "✅  AFTER TravAll\n\n• Instant Green Passport alerts\n• AI adjusts schedule in seconds\n• Single Super App replaces 3-4\n• 20% time saved → pays subscription",
         imgs[1]),
        ("Scenario 3: Melih Özgen", ORANGE,
         "❌  BEFORE TravAll\n\n• First in family abroad — stressed\n• Visa so confusing → may cancel\n• Frustrated by complex tables\n• Switches off phone, gives up",
         "✅  AFTER TravAll\n\n• Smart Assistant greets calmly\n• Visual icons: 90% success routes\n• Uploads docs with confirmations\n• Final Approval = full control",
         imgs[2]),
    ]

    for title, accent, before_text, after_text, img_path in scenarios:
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(sl, WHITE)

        # Title
        tb = sl.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = DARK_BLUE

        # Before/After as TABLE (2 rows, 1 col) on the left
        table_shape = sl.shapes.add_table(2, 1, Inches(0.2), Inches(0.8), Inches(4.4), Inches(6.3))
        table = table_shape.table
        no_border(table)
        table.columns[0].width = Inches(4.4)

        set_cell(table.cell(0, 0), before_text, 13, False, G700, PP_ALIGN.LEFT, LIGHT_RED)
        # Make header bold
        tf = table.cell(0, 0).text_frame
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RED
        tf.paragraphs[0].font.size = Pt(16)

        set_cell(table.cell(1, 0), after_text, 13, False, G700, PP_ALIGN.LEFT, LIGHT_GREEN)
        tf = table.cell(1, 0).text_frame
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = GREEN
        tf.paragraphs[0].font.size = Pt(16)

        # Storyboard image on the right
        if os.path.exists(img_path):
            sl.shapes.add_picture(img_path, Inches(4.85), Inches(0.6), Inches(8.2), Inches(6.6))

    # ===== SLIDE 7: RESULTS & TAKEAWAYS =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, WHITE)

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Design Requirements for Phase 4"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = DARK_BLUE

    # Takeaways as TABLE (5 rows, 2 cols)
    table_shape = sl.shapes.add_table(5, 2, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.8))
    table = table_shape.table
    no_border(table)
    table.columns[0].width = Inches(5.0)
    table.columns[1].width = Inches(7.3)

    takeaways = [
        ("📱  Mobile-First & One-Handed", "62.26% use phone one-handed → buttons at bottom of screen", MED_BLUE),
        ("👁️  Visual Hierarchy & Summarization", "Dashboard with icons & brief summaries (<140 chars) — no dense tables", TEAL),
        ("🔐  User Control & Transparency", "94.34% want final approval themselves — AI assists but never decides", GREEN),
        ("✨  Minimalist Interface (Anti-Muda)", "73.59% prefer minimalist UI over feature overload — simplest flows", ORANGE),
        ("🔔  Dynamic Notifications & Live Data", "86.79% want proactive border/rule change alerts for faster decisions", RED),
    ]

    for i, (title, desc, c) in enumerate(takeaways):
        set_cell(table.cell(i, 0), title, 17, True, DARK_BLUE, PP_ALIGN.LEFT, G100)
        set_cell(table.cell(i, 1), desc, 15, False, G500, PP_ALIGN.LEFT, G100)

    # ===== SLIDE 8: THANK YOU =====
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(sl, DARK_BLUE)

    tb = sl.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(52); p.font.bold = True; p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(30)

    p = tf.add_paragraph()
    p.text = "Questions & Discussion"
    p.font.size = Pt(26); p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(40)

    p = tf.add_paragraph()
    p.text = "Özde  •  Nilsu  •  Deniz  •  Fatih Bilal"
    p.font.size = Pt(18); p.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    p.alignment = PP_ALIGN.CENTER

    out = os.path.join(OUT_DIR, "IE48L_Phase3_Presentation.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    return out


if __name__ == '__main__':
    print("Extracting storyboards...")
    images = extract_images(PDF_PATH, OUT_DIR)
    print("Building presentation...")
    build(images)
    print("Done!")
