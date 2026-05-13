#!/usr/bin/env python3
"""Generate TRAVALL Phase 5B & 6B Presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(os.path.abspath(__file__))

# Brand colors
DARK_BLUE = RGBColor(0x0A, 0x1F, 0x3C)
MID_BLUE = RGBColor(0x14, 0x3D, 0x6B)
ACCENT_BLUE = RGBColor(0x1E, 0x88, 0xE5)
LIGHT_BLUE = RGBColor(0xBB, 0xDE, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)
GOLD = RGBColor(0xFF, 0xB3, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_gradient_bg(slide, color1=DARK_BLUE, color2=MID_BLUE):
    """Add a solid dark background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color1

def add_shape_rect(slide, left, top, width, height, fill_color, alpha=None):
    """Add a colored rectangle"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=WHITE, bold=False, spacing=1.2, font_name="Calibri"):
    """Add text box with multiple paragraphs"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.5)
    return txBox

def add_icon_text_pair(slide, left, top, icon_text, label, value, label_size=14, value_size=28):
    """Add an icon-style metric card"""
    card = add_rounded_rect(slide, left, top, Inches(2.3), Inches(1.4), MID_BLUE)
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.2), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].text = icon_text
    ctf.paragraphs[0].font.size = Pt(16)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.word_wrap = False

    add_text_box(slide, left + Inches(0.8), top + Inches(0.15), Inches(1.4), Inches(0.35),
                 value, font_size=value_size, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.8), top + Inches(0.75), Inches(1.4), Inches(0.5),
                 label, font_size=label_size, color=LIGHT_BLUE)
    return card

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_gradient_bg(slide)

# Decorative accent bar at top
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

# Left side - text content
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(3), Inches(0.5),
             "IE48L", font_size=16, color=LIGHT_BLUE, bold=False)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(1.2),
             "TRAVALL", font_size=60, color=WHITE, bold=True, font_name="Calibri Light")

add_text_box(slide, Inches(0.8), Inches(2.8), Inches(7), Inches(0.6),
             "Prototype & Field Trial-2", font_size=28, color=ACCENT_BLUE, bold=False, font_name="Calibri Light")

add_text_box(slide, Inches(0.8), Inches(3.6), Inches(7), Inches(0.4),
             "Phase 5B & 6B", font_size=20, color=LIGHT_BLUE)

# Decorative line
add_shape_rect(slide, Inches(0.8), Inches(4.2), Inches(4), Inches(0.04), ACCENT_BLUE)

# Team members
team_lines = [
    "Ozde Avdan  -  2021402126",
    "Nilsu Calliogullari  -  2021402165",
    "Deniz Tas  -  2021402123",
    "Fatih Bilal Yilmaz  -  2021402174"
]
add_multi_text(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(2),
               team_lines, font_size=14, color=LIGHT_BLUE, spacing=1.5)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(3), Inches(0.4),
             "May 4, 2026", font_size=14, color=RGBColor(0x90, 0xA4, 0xAE))

# Right side - Logo
logo_path = os.path.join(BASE, "bogazici_logo.png")
if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(10.5), Inches(0.5), Inches(1.8), Inches(1.8))

# Right side - hi-fi mockup preview
hifi_img = os.path.join(BASE, "img_p11_0.png")
if os.path.exists(hifi_img):
    slide.shapes.add_picture(hifi_img, Inches(8.5), Inches(2.5), Inches(4.2), Inches(4.5))

# ============================================================
# SLIDE 2: AGENDA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
             "Agenda", font_size=40, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(2.5), Inches(0.04), ACCENT_BLUE)

agenda_items = [
    ("01", "Field Trial-1 Recap", "Key findings from Lo-Fi testing"),
    ("02", "Usability Problems", "3 critical issues identified"),
    ("03", "Design Improvements", "Lo-Fi to Hi-Fi evolution"),
    ("04", "Hi-Fi Prototype", "Screen showcase"),
    ("05", "Field Trial-2 Results", "31 participants, 15 tasks"),
    ("06", "Comparison & Feedback", "Lo-Fi vs Hi-Fi metrics"),
    ("07", "Future Roadmap", "Next steps for Travall"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.8) + Inches(i * 0.72)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].text = num
    ctf.paragraphs[0].font.size = Pt(16)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.word_wrap = False

    add_text_box(slide, Inches(2.0), y - Inches(0.02), Inches(4), Inches(0.35),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.32), Inches(4), Inches(0.3),
                 desc, font_size=13, color=LIGHT_BLUE)

# Right side - decorative app screenshots
right_img = os.path.join(BASE, "img_p14_0.png")
if os.path.exists(right_img):
    slide.shapes.add_picture(right_img, Inches(8), Inches(1), Inches(4.5), Inches(5.5))

# ============================================================
# SLIDE 3: FIELD TRIAL-1 RECAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
             "Field Trial-1 Recap (Lo-Fi)", font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(0.5),
             "5 participants  |  15 sequential tasks  |  Lo-Fi prototype",
             font_size=16, color=LIGHT_BLUE)

# Metric cards
metrics_y = Inches(2.3)
add_icon_text_pair(slide, Inches(0.8), metrics_y, "T", "Avg. Completion", "6'54\"")
add_icon_text_pair(slide, Inches(3.4), metrics_y, "%", "Success Rate", "96%")
add_icon_text_pair(slide, Inches(6.0), metrics_y, "E", "Avg. Errors", "2.2")

# Subjective ratings
ratings_y = Inches(4.0)
add_text_box(slide, Inches(0.8), ratings_y, Inches(5), Inches(0.4),
             "Subjective Ratings", font_size=20, color=WHITE, bold=True)

rating_items = [
    ("Visual Clarity", "9.0/10", GREEN),
    ("Likability", "8.6/10", ACCENT_BLUE),
    ("Nav. Comfort", "8.2/10", ACCENT_BLUE),
    ("Ease of Use", "7.8/10", ORANGE),
]
for i, (label, val, col) in enumerate(rating_items):
    ry = ratings_y + Inches(0.55) + Inches(i * 0.55)
    add_text_box(slide, Inches(1.0), ry, Inches(2.5), Inches(0.4),
                 label, font_size=15, color=LIGHT_BLUE)
    # Bar background
    add_rounded_rect(slide, Inches(3.2), ry + Inches(0.05), Inches(3.5), Inches(0.3), RGBColor(0x1A, 0x2E, 0x4A))
    # Bar fill
    score = float(val.split("/")[0])
    bar_w = Inches(3.5 * score / 10.0)
    add_rounded_rect(slide, Inches(3.2), ry + Inches(0.05), bar_w, Inches(0.3), col)
    add_text_box(slide, Inches(6.85), ry, Inches(1), Inches(0.4),
                 val, font_size=14, color=WHITE, bold=True)

# Right side - lo-fi aggregate table image
lofi_agg = os.path.join(BASE, "lofi_p4_1.png")
if os.path.exists(lofi_agg):
    slide.shapes.add_picture(lofi_agg, Inches(8.5), Inches(2), Inches(4), Inches(3))

# ============================================================
# SLIDE 4: 3 USABILITY PROBLEMS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "3 Usability Problems Identified", font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT_BLUE)

problems = [
    ("1", "Settings Discoverability",
     "Users couldn't find notification settings.\nAccount section blended with content.\nExpected a gear icon or dedicated tab.",
     RED),
    ("2", "Apply vs. Renew Confusion",
     "Two CTA buttons looked identical.\nUsers tapped wrong button.\nSpent time on the wrong screen.",
     ORANGE),
    ("3", "Dashed Border = Inactive",
     "Add Second Passport used dashed border.\nPerceived as disabled placeholder.\nNot recognized as interactive.",
     GOLD),
]

card_width = Inches(3.6)
for i, (num, title, desc, accent) in enumerate(problems):
    x = Inches(0.8) + Inches(i * 3.9)
    y = Inches(1.8)

    card = add_rounded_rect(slide, x, y, card_width, Inches(3.8), MID_BLUE)

    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.2), Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    btf = badge.text_frame
    btf.paragraphs[0].text = num
    btf.paragraphs[0].font.size = Pt(22)
    btf.paragraphs[0].font.color.rgb = WHITE
    btf.paragraphs[0].font.bold = True
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.85), y + Inches(0.25), Inches(2.5), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True)

    # Accent bar
    add_shape_rect(slide, x + Inches(0.15), y + Inches(0.95), Inches(3.3), Inches(0.03), accent)

    lines = desc.split("\n")
    for j, line in enumerate(lines):
        add_text_box(slide, x + Inches(0.25), y + Inches(1.2) + Inches(j * 0.5), Inches(3.1), Inches(0.45),
                     line, font_size=14, color=LIGHT_BLUE)

# Lo-fi screenshot showing the problematic dashed border
lofi_passport = os.path.join(BASE, "lofi_p13_0.png")
if os.path.exists(lofi_passport):
    slide.shapes.add_picture(lofi_passport, Inches(8.5), Inches(5.7), Inches(2.2), Inches(1.7))
    add_text_box(slide, Inches(8.5), Inches(5.45), Inches(2.2), Inches(0.3),
                 "Lo-Fi: Dashed Border", font_size=11, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Heuristic evaluation table
heur_img = os.path.join(BASE, "img_p3_0.png")
if os.path.exists(heur_img):
    slide.shapes.add_picture(heur_img, Inches(0.5), Inches(5.8), Inches(7.5), Inches(1.5))
    add_text_box(slide, Inches(0.5), Inches(5.55), Inches(7.5), Inches(0.3),
                 "Nielsen's 10 Heuristics Assessment", font_size=11, color=ACCENT_BLUE, bold=True)

# ============================================================
# SLIDE 5: DESIGN FIXES - BEFORE / AFTER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Design Fixes: Lo-Fi  -->  Hi-Fi", font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_BLUE)

# Fix descriptions on the left
fixes = [
    ("Fix 1: Settings Access",
     "Added bell icon to header\n1 tap instead of 4 taps"),
    ("Fix 2: Button Differentiation",
     "Apply = outlined (secondary)\nRenew = filled (primary)"),
    ("Fix 3: Solid Button",
     "Dashed border replaced\nwith standard CTA styling"),
]

for i, (title, desc) in enumerate(fixes):
    y = Inches(1.4) + Inches(i * 1.6)
    # Accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = GREEN
    dot.line.fill.background()

    add_text_box(slide, Inches(1.2), y, Inches(3.5), Inches(0.35),
                 title, font_size=17, color=WHITE, bold=True)
    lines = desc.split("\n")
    for j, line in enumerate(lines):
        add_text_box(slide, Inches(1.2), y + Inches(0.35) + Inches(j * 0.35), Inches(3.5), Inches(0.35),
                     line, font_size=13, color=LIGHT_BLUE)

# Lo-fi vs Hi-fi screenshots side by side
add_text_box(slide, Inches(5.0), Inches(1.2), Inches(3.5), Inches(0.35),
             "Lo-Fi (Before)", font_size=16, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.0), Inches(1.2), Inches(3.5), Inches(0.35),
             "Hi-Fi (After)", font_size=16, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Lo-fi Passport Hub
lofi_hub = os.path.join(BASE, "lofi_p9_0.png")
if os.path.exists(lofi_hub):
    slide.shapes.add_picture(lofi_hub, Inches(5.0), Inches(1.6), Inches(3.5), Inches(5.5))

# Hi-fi Passport Hub
hifi_hub = os.path.join(BASE, "img_p11_0.png")
if os.path.exists(hifi_hub):
    slide.shapes.add_picture(hifi_hub, Inches(9.0), Inches(1.6), Inches(3.5), Inches(5.5))

# ============================================================
# SLIDE 6: HI-FI PROTOTYPE - PASSPORT & VISA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Hi-Fi Prototype: Passport & Visa Hub", font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_BLUE)

# Screenshots
img1 = os.path.join(BASE, "img_p11_0.png")
img2 = os.path.join(BASE, "img_p12_0.png")
if os.path.exists(img1):
    slide.shapes.add_picture(img1, Inches(0.5), Inches(1.4), Inches(6), Inches(5.8))
if os.path.exists(img2):
    slide.shapes.add_picture(img2, Inches(6.8), Inches(1.4), Inches(6), Inches(5.8))

# ============================================================
# SLIDE 7: HI-FI PROTOTYPE - TRIPS & EXPLORE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Hi-Fi Prototype: Trips & Explore", font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_BLUE)

img3 = os.path.join(BASE, "img_p13_0.png")
img4 = os.path.join(BASE, "img_p14_0.png")
if os.path.exists(img3):
    slide.shapes.add_picture(img3, Inches(0.5), Inches(1.4), Inches(6), Inches(5.8))
if os.path.exists(img4):
    slide.shapes.add_picture(img4, Inches(6.8), Inches(1.4), Inches(6), Inches(5.8))

# ============================================================
# SLIDE 8: HI-FI PROTOTYPE - PROFILE & SETTINGS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Hi-Fi Prototype: Profile & Settings", font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_BLUE)

img5 = os.path.join(BASE, "img_p15_0.png")
if os.path.exists(img5):
    slide.shapes.add_picture(img5, Inches(1.5), Inches(1.4), Inches(5.5), Inches(5.8))

# System map
sys_map = os.path.join(BASE, "img_p10_0.png")
if os.path.exists(sys_map):
    slide.shapes.add_picture(sys_map, Inches(7.5), Inches(1.2), Inches(5), Inches(6))
    add_text_box(slide, Inches(7.5), Inches(6.9), Inches(5), Inches(0.35),
                 "Revised System Map", font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: FIELD TRIAL-2 RESULTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Field Trial-2 Results", font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(8), Inches(0.5),
             "31 participants  |  15 sequential tasks  |  Hi-Fi interactive prototype",
             font_size=16, color=LIGHT_BLUE)

# Big metric cards
metrics = [
    ("4'58\"", "Avg. Time", ACCENT_BLUE),
    ("98.3%", "Success Rate", GREEN),
    ("0.8", "Avg. Errors", GREEN),
    ("9.2", "Likability", ACCENT_BLUE),
    ("8.8", "Ease of Use", ACCENT_BLUE),
    ("9.5", "Visual Clarity", GREEN),
]

for i, (val, label, accent) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.0)
    y = Inches(2.0) + Inches(row * 2.2)

    card = add_rounded_rect(slide, x, y, Inches(3.5), Inches(1.8), MID_BLUE)

    # Accent bar at top of card
    add_shape_rect(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.2), Inches(0.05), accent)

    add_text_box(slide, x, y + Inches(0.3), Inches(3.5), Inches(0.9),
                 val, font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri Light")
    add_text_box(slide, x, y + Inches(1.2), Inches(3.5), Inches(0.4),
                 label, font_size=16, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# Navigation comfort
add_rounded_rect(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8), MID_BLUE)
add_text_box(slide, Inches(1.0), Inches(6.5), Inches(5.5), Inches(0.6),
             "Navigation Comfort: 9.2/10", font_size=20, color=WHITE, bold=True)
add_text_box(slide, Inches(6.5), Inches(6.5), Inches(5.5), Inches(0.6),
             "Bell icon used by 22/31 participants for notification access",
             font_size=15, color=LIGHT_BLUE)

# ============================================================
# SLIDE 10: COMPARISON LO-FI vs HI-FI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Lo-Fi vs Hi-Fi: Side-by-Side", font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_BLUE)

# Comparison table
comp_img = os.path.join(BASE, "img_p7_1.png")
if os.path.exists(comp_img):
    slide.shapes.add_picture(comp_img, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3))

# Visual comparison bars
comparisons = [
    ("Completion Time", "6'54\"", "4'58\"", 6.9, 4.97, "28% faster"),
    ("Success Rate", "96%", "98.3%", 96, 98.3, "+2.3%"),
    ("Ease of Use", "7.8", "8.8", 78, 88, "+12.8%"),
    ("Likability", "8.6", "9.2", 86, 92, "+7%"),
    ("Visual Clarity", "9.0", "9.5", 90, 95, "+5.6%"),
    ("Nav. Comfort", "8.2", "9.2", 82, 92, "+12.2%"),
]

start_x = Inches(7)
bar_max = Inches(4.5)

add_text_box(slide, Inches(7), Inches(1.5), Inches(2), Inches(0.3),
             "Lo-Fi", font_size=13, color=ORANGE, bold=True)
add_text_box(slide, Inches(8.5), Inches(1.5), Inches(2), Inches(0.3),
             "Hi-Fi", font_size=13, color=GREEN, bold=True)

for i, (label, v1, v2, p1, p2, delta) in enumerate(comparisons):
    y = Inches(2.0) + Inches(i * 0.85)

    add_text_box(slide, Inches(7), y, Inches(2.5), Inches(0.3),
                 label, font_size=14, color=LIGHT_BLUE)

    # Lo-fi bar
    add_rounded_rect(slide, Inches(9.2), y, Inches(2.8 * min(p1, 100) / 100), Inches(0.25), ORANGE)
    add_text_box(slide, Inches(9.2) + Inches(2.8 * min(p1, 100) / 100) + Inches(0.05), y - Inches(0.02),
                 Inches(0.8), Inches(0.3), v1, font_size=11, color=ORANGE, bold=True)

    # Hi-fi bar
    add_rounded_rect(slide, Inches(9.2), y + Inches(0.28), Inches(2.8 * min(p2, 100) / 100), Inches(0.25), GREEN)
    add_text_box(slide, Inches(9.2) + Inches(2.8 * min(p2, 100) / 100) + Inches(0.05), y + Inches(0.26),
                 Inches(0.8), Inches(0.3), v2, font_size=11, color=GREEN, bold=True)

    # Delta badge
    add_text_box(slide, Inches(12.2), y + Inches(0.1), Inches(1), Inches(0.3),
                 delta, font_size=12, color=GREEN, bold=True)

# Key insight box
add_rounded_rect(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.2), MID_BLUE)
add_text_box(slide, Inches(1.2), Inches(5.2), Inches(10.5), Inches(0.4),
             "Key Findings", font_size=20, color=WHITE, bold=True)
insights = [
    "All 3 usability problems from Trial-1 were fully resolved",
    "Apply/Renew button confusion: 0 errors (was the #1 issue)",
    "Notification settings: 100% success rate with bell icon shortcut",
    "Add Second Passport: 31/31 tapped correctly (dashed border eliminated)",
]
for i, insight in enumerate(insights):
    add_text_box(slide, Inches(1.5), Inches(5.7) + Inches(i * 0.38), Inches(10), Inches(0.35),
                 "   " + insight, font_size=14, color=LIGHT_BLUE)

# ============================================================
# SLIDE 11: PARTICIPANT FEEDBACK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Participant Feedback", font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_BLUE)

# Positive feedback - left column
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(0.4),
             "What They Loved", font_size=18, color=GREEN, bold=True)

positive = [
    "\"It looks like a real app, not a prototype.\"",
    "\"The bell button is exactly where I expected it.\"",
    "\"Apply and Renew buttons are clearly different.\"",
    "\"The donut chart is a nice touch.\"",
    "\"Timeline view is perfectly structured.\"",
    "\"Community feed budget breakdowns are invaluable.\"",
]

for i, quote in enumerate(positive):
    y = Inches(2.0) + Inches(i * 0.8)
    card = add_rounded_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.65), MID_BLUE)
    # Green left accent
    add_shape_rect(slide, Inches(0.8), y, Inches(0.06), Inches(0.65), GREEN)
    add_text_box(slide, Inches(1.1), y + Inches(0.08), Inches(5.0), Inches(0.5),
                 quote, font_size=13, color=LIGHT_BLUE)

# Negative feedback - right column
add_text_box(slide, Inches(7), Inches(1.5), Inches(3), Inches(0.4),
             "Improvement Ideas", font_size=18, color=ORANGE, bold=True)

negative = [
    "\"World map should be interactive.\"",
    "\"Search should filter as I type.\"",
    "\"No form validation on visa application.\"",
    "\"Would like to compare countries.\"",
    "\"Dark mode toggle doesn't work yet.\"",
    "\"Need search/filter in Explore feed.\"",
]

for i, quote in enumerate(negative):
    y = Inches(2.0) + Inches(i * 0.8)
    card = add_rounded_rect(slide, Inches(7), y, Inches(5.5), Inches(0.65), MID_BLUE)
    add_shape_rect(slide, Inches(7), y, Inches(0.06), Inches(0.65), ORANGE)
    add_text_box(slide, Inches(7.3), y + Inches(0.08), Inches(5.0), Inches(0.5),
                 quote, font_size=13, color=LIGHT_BLUE)

# ============================================================
# SLIDE 12: FUTURE ROADMAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Future Roadmap", font_size=36, color=WHITE, bold=True, font_name="Calibri Light")
add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_BLUE)

roadmap_items = [
    ("Interactive World Map", "Real SVG map with visited countries colored", "MAP"),
    ("Live Search", "Real-time filtering on Passport Hub country list", "SRC"),
    ("Form Validation", "Input validation & error states for visa forms", "VAL"),
    ("Country Comparison", "Compare visa requirements of 2-3 countries", "CMP"),
    ("Community Filters", "Destination & budget filters for Explore feed", "FLT"),
    ("Dark Mode", "Full dark mode integration with state persistence", "DRK"),
    ("Trip Sharing", "Share trip details via OS share sheet", "SHR"),
    ("User Onboarding", "First-time guide highlighting key features", "ONB"),
]

for i, (title, desc, icon) in enumerate(roadmap_items):
    col = i % 4
    row = i // 4
    x = Inches(0.6) + Inches(col * 3.15)
    y = Inches(1.6) + Inches(row * 2.8)

    card = add_rounded_rect(slide, x, y, Inches(2.9), Inches(2.3), MID_BLUE)

    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.05), y + Inches(0.2), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].text = icon
    ctf.paragraphs[0].font.size = Pt(14)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.word_wrap = False

    add_text_box(slide, x + Inches(0.1), y + Inches(1.05), Inches(2.7), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.45), Inches(2.6), Inches(0.7),
                 desc, font_size=12, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13: THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

# Big centered text
add_text_box(slide, Inches(0), Inches(1.5), Inches(13.333), Inches(1.5),
             "Thank You!", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

add_text_box(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(0.6),
             "TRAVALL  -  Your Smart Travel Companion", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(0), Inches(4.2), Inches(13.333), Inches(0.5),
             "Questions?", font_size=28, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

# Team members centered
team_text = "Ozde Avdan  |  Nilsu Calliogullari  |  Deniz Tas  |  Fatih Bilal Yilmaz"
add_text_box(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(0.5),
             team_text, font_size=16, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(6.1), Inches(13.333), Inches(0.4),
             "Bogazici University  |  IE48L  |  May 2026", font_size=14,
             color=RGBColor(0x90, 0xA4, 0xAE), alignment=PP_ALIGN.CENTER)

# Logo centered
if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(6.0), Inches(6.5), Inches(0.9), Inches(0.9))

# Bottom accent bar
add_shape_rect(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT_BLUE)

# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(BASE, "TRAVALL_Phase5B_6B_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
