#!/usr/bin/env python3
"""
Generate Phase 4 PPTX presentation for TravAll - IE48L
5-minute presentation with app screenshots
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PHASE_DIR = "/Users/fatihbilalyilmaz/PycharmProjects/IE492/IE48L/Phase 4"
OUT = os.path.join(PHASE_DIR, "IE48L-Phase4-Presentation.pptx")
LOGO = os.path.join(PHASE_DIR, "extracted_img_0.png")

# Screenshots
SS = [os.path.join(PHASE_DIR, f"Ekran Resmi 2026-03-30 12.14.{t}.png")
      for t in ["18", "23", "32", "39", "46", "53", "59"]]

# Diagrams (cropped from PDF)
SYSMAP = os.path.join(PHASE_DIR, "system_map_cropped.png")
FLOW_DASH = os.path.join(PHASE_DIR, "flow_dashboard.png")
FLOW_VISA = os.path.join(PHASE_DIR, "flow_visa.png")
FLOW_JOURNEY = os.path.join(PHASE_DIR, "flow_journey.png")
FLOW_COMMUNITY = os.path.join(PHASE_DIR, "flow_community.png")
FLOW_ANALYTICS = os.path.join(PHASE_DIR, "flow_analytics_cropped.png")

# Colors
DARK_BLUE = RGBColor(0x0D, 0x2B, 0x45)
ACCENT_BLUE = RGBColor(0x1A, 0x73, 0xE8)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x7B, 0x8D)
MED_GRAY = RGBColor(0x95, 0xA5, 0xB5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BLUE):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bold_first=False, spacing=Pt(6)):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def add_footer(slide, slide_num, total=9):
    """Add footer bar."""
    bar = add_shape(slide, Inches(0), H - Inches(0.45), W, Inches(0.45),
                    RGBColor(0x08, 0x1F, 0x35))
    add_text(slide, Inches(0.5), H - Inches(0.38), Inches(4), Inches(0.3),
             "IE48L  |  TravAll  |  Phase 4", font_size=10, color=MED_GRAY)
    add_text(slide, W - Inches(1.5), H - Inches(0.38), Inches(1.2), Inches(0.3),
             f"{slide_num} / {total}", font_size=10, color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# Left accent bar
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

# Logo
if os.path.exists(LOGO):
    slide.shapes.add_picture(LOGO, Inches(5.7), Inches(0.6), Inches(1.8), Inches(1.8))

# Title
add_text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.7),
         "Phase 4", font_size=48, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.9),
         "Task Analysis, System Map and Screen Designs",
         font_size=28, color=WHITE, bold=True)

# Course
add_text(slide, Inches(1.5), Inches(4.5), Inches(5), Inches(0.5),
         "IE48L  -  Human Computer Interaction Design", font_size=16, color=MED_GRAY)

# Date
add_text(slide, Inches(1.5), Inches(5.0), Inches(5), Inches(0.4),
         "30.03.2026", font_size=14, color=MED_GRAY)

# Team members - right side
members = [
    "Ozde Avdan - 2021402126",
    "Nilsu Calliogullari - 2021402165",
    "Deniz Tas - 2021402123",
    "Fatih Bilal Yilmaz - 2021402174"
]
add_bullet_text(slide, Inches(8.5), Inches(5.0), Inches(4), Inches(2),
                members, font_size=13, color=MED_GRAY)

add_footer(slide, 1, total=11)

# ============================================================
# SLIDE 2: AGENDA / OVERVIEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.7),
         "Agenda", font_size=36, color=WHITE, bold=True)

# Divider line
line = add_shape(slide, Inches(1), Inches(1.2), Inches(3), Inches(0.04), ACCENT_BLUE)

sections = [
    ("01", "Task Analysis", "3 user tasks mapped to personas from Phase 3"),
    ("02", "System Map", "Complete navigation hierarchy of TravAll"),
    ("03", "Screen List", "Detailed UI element documentation per screen"),
    ("04", "Screen Designs", "High-fidelity prototypes for all core modules"),
]

y = Inches(1.8)
for num, title, desc in sections:
    # Number circle
    circ = add_rounded_rect(slide, Inches(1.2), y, Inches(0.8), Inches(0.8), ACCENT_BLUE)
    circ.text_frame.paragraphs[0].text = num
    circ.text_frame.paragraphs[0].font.size = Pt(22)
    circ.text_frame.paragraphs[0].font.color.rgb = WHITE
    circ.text_frame.paragraphs[0].font.bold = True
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circ.text_frame.paragraphs[0].font.name = "Calibri"

    add_text(slide, Inches(2.3), y + Inches(0.05), Inches(4), Inches(0.4),
             title, font_size=22, color=WHITE, bold=True)
    add_text(slide, Inches(2.3), y + Inches(0.42), Inches(5), Inches(0.35),
             desc, font_size=14, color=MED_GRAY)
    y += Inches(1.2)

# Right side: quick app preview
if os.path.exists(SS[0]):
    slide.shapes.add_picture(SS[0], Inches(7.8), Inches(0.5), Inches(5), Inches(6.5))

add_footer(slide, 2, total=11)

# ============================================================
# SLIDE 3: TASK ANALYSIS - Task 1
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.4), Inches(8), Inches(0.6),
         "1. Task Analysis", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(1.0), Inches(2.5), Inches(0.04), ACCENT_BLUE)

# Task 1 box
t1_box = add_rounded_rect(slide, Inches(0.6), Inches(1.3), Inches(6), Inches(1.0),
                           RGBColor(0x14, 0x3A, 0x5C))
add_text(slide, Inches(0.8), Inches(1.35), Inches(5.5), Inches(0.35),
         "Task 1: Visa Status Check", font_size=18, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.55),
         "A Green Passport holder plans a trip to Germany - confirm visa-free status and check travel advisories.",
         font_size=12, color=MED_GRAY)

# Table for Task 1
steps_t1 = [
    ("1. Authenticate Profile", "Input passport type during enrollment", "System filters by entitlement"),
    ("2. Access Visa Assistant", "Click Dynamic Visa Assistant button", "Search bar for destinations"),
    ("3. Search Destination", "Type 'Germany' in search bar", "Auto-sifts for Green Passports"),
    ("4. Review Summary", "Read icons and short summaries", "'Visa Free' status with details"),
]
y = Inches(2.5)
# Header
hdr = add_rounded_rect(slide, Inches(0.6), y, Inches(6), Inches(0.4), ACCENT_BLUE)
add_text(slide, Inches(0.7), y + Inches(0.05), Inches(1.8), Inches(0.3),
         "Step", font_size=11, color=WHITE, bold=True)
add_text(slide, Inches(2.5), y + Inches(0.05), Inches(1.8), Inches(0.3),
         "Way Performed", font_size=11, color=WHITE, bold=True)
add_text(slide, Inches(4.4), y + Inches(0.05), Inches(2.1), Inches(0.3),
         "Feedback", font_size=11, color=WHITE, bold=True)

y += Inches(0.45)
for step, way, feedback in steps_t1:
    row = add_rounded_rect(slide, Inches(0.6), y, Inches(6), Inches(0.42),
                            RGBColor(0x14, 0x3A, 0x5C))
    add_text(slide, Inches(0.7), y + Inches(0.06), Inches(1.7), Inches(0.32),
             step, font_size=10, color=WHITE)
    add_text(slide, Inches(2.5), y + Inches(0.06), Inches(1.8), Inches(0.32),
             way, font_size=10, color=MED_GRAY)
    add_text(slide, Inches(4.4), y + Inches(0.06), Inches(2.1), Inches(0.32),
             feedback, font_size=10, color=MED_GRAY)
    y += Inches(0.45)

# Task 2 box
t2_box = add_rounded_rect(slide, Inches(0.6), y + Inches(0.15), Inches(6), Inches(0.85),
                           RGBColor(0x14, 0x3A, 0x5C))
add_text(slide, Inches(0.8), y + Inches(0.2), Inches(5.5), Inches(0.3),
         "Task 2: Museum Ticket & Directions", font_size=18, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), y + Inches(0.52), Inches(5.5), Inches(0.4),
         "Find nearby museum, purchase digital ticket, get directions. (4 steps: Explore > Purchase > Access > Navigate)",
         font_size=11, color=MED_GRAY)

# Task 3 box
y2 = y + Inches(1.15)
t3_box = add_rounded_rect(slide, Inches(0.6), y2, Inches(6), Inches(0.85),
                           RGBColor(0x14, 0x3A, 0x5C))
add_text(slide, Inches(0.8), y2 + Inches(0.05), Inches(5.5), Inches(0.3),
         "Task 3: Blog & Statistics Review", font_size=18, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), y2 + Inches(0.37), Inches(5.5), Inches(0.4),
         "Erasmus student shares visa experience on blog, reviews spending % and visited country progress. (5 steps)",
         font_size=11, color=MED_GRAY)

# Right: Visa screenshot
if os.path.exists(SS[2]):
    slide.shapes.add_picture(SS[2], Inches(7.2), Inches(0.5), Inches(5.5), Inches(6.5))

add_footer(slide, 3, total=11)

# ============================================================
# SLIDE 4: SYSTEM MAP (Full Diagram)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

# Dark header bar
add_shape(slide, Inches(0.15), Inches(0), W, Inches(1.0), DARK_BLUE)
add_text(slide, Inches(1), Inches(0.2), Inches(8), Inches(0.6),
         "2. System Map - Full Navigation Architecture", font_size=32, color=WHITE, bold=True)

# System map image - centered, large, on white background
if os.path.exists(SYSMAP):
    # Image is 1725x1670 -> ratio ~1.03:1
    img_w = Inches(8.5)
    img_h = Inches(5.6)
    left = (W - img_w) // 2
    slide.shapes.add_picture(SYSMAP, left, Inches(1.2), img_w, img_h)

# Bottom label
add_rounded_rect(slide, Inches(3.5), Inches(6.9), Inches(6.3), Inches(0.4), ACCENT_BLUE)
add_text(slide, Inches(3.6), Inches(6.92), Inches(6.1), Inches(0.35),
         "5-Tab Navigation  |  Max 3-Level Depth  |  30+ Screens  |  Cross-Module Links",
         font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_footer(slide, 4, total=11)

# ============================================================
# SLIDE 5: USER FLOWS - Dashboard & Visa
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_shape(slide, Inches(0.15), Inches(0), W, Inches(1.0), DARK_BLUE)
add_text(slide, Inches(1), Inches(0.2), Inches(10), Inches(0.6),
         "2b. User Flows - Smart Dashboard & Visa Assistant", font_size=30, color=WHITE, bold=True)

# Dashboard flow - left side
if os.path.exists(FLOW_DASH):
    slide.shapes.add_picture(FLOW_DASH, Inches(0.3), Inches(1.15), Inches(7.5), Inches(3.6))

# Label
add_rounded_rect(slide, Inches(0.5), Inches(4.85), Inches(3.0), Inches(0.35), ACCENT_BLUE)
add_text(slide, Inches(0.6), Inches(4.87), Inches(2.8), Inches(0.3),
         "Smart Dashboard Mapping", font_size=11, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

# Visa flow - right side
if os.path.exists(FLOW_VISA):
    # Visa diagram is shorter (670px high)
    slide.shapes.add_picture(FLOW_VISA, Inches(3.8), Inches(5.3), Inches(4.5), Inches(1.7))

add_rounded_rect(slide, Inches(4.0), Inches(7.05), Inches(2.8), Inches(0.35), ACCENT_BLUE)
# removed - goes below footer

# Right column: Journey flow preview
if os.path.exists(FLOW_JOURNEY):
    slide.shapes.add_picture(FLOW_JOURNEY, Inches(8.3), Inches(1.15), Inches(4.7), Inches(4.4))

add_rounded_rect(slide, Inches(8.8), Inches(5.65), Inches(3.5), Inches(0.35), ACCENT_BLUE)
add_text(slide, Inches(8.9), Inches(5.67), Inches(3.3), Inches(0.3),
         "Journey Timeline Mapping", font_size=11, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_footer(slide, 5, total=11)

# ============================================================
# SLIDE 6: USER FLOWS - Community & Analytics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_shape(slide, Inches(0.15), Inches(0), W, Inches(1.0), DARK_BLUE)
add_text(slide, Inches(1), Inches(0.2), Inches(10), Inches(0.6),
         "2c. User Flows - Community Hub & Travel Analytics", font_size=30, color=WHITE, bold=True)

# Community flow - left
if os.path.exists(FLOW_COMMUNITY):
    slide.shapes.add_picture(FLOW_COMMUNITY, Inches(0.3), Inches(1.2), Inches(6.5), Inches(4.2))

add_rounded_rect(slide, Inches(1.0), Inches(5.5), Inches(3.2), Inches(0.35), ACCENT_BLUE)
add_text(slide, Inches(1.1), Inches(5.52), Inches(3.0), Inches(0.3),
         "Community Hub Mapping", font_size=11, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

# Analytics flow - right
if os.path.exists(FLOW_ANALYTICS):
    slide.shapes.add_picture(FLOW_ANALYTICS, Inches(7.0), Inches(1.2), Inches(5.8), Inches(3.5))

add_rounded_rect(slide, Inches(8.0), Inches(4.85), Inches(3.5), Inches(0.35), ACCENT_BLUE)
add_text(slide, Inches(8.1), Inches(4.87), Inches(3.3), Inches(0.3),
         "Travel Analytics Mapping", font_size=11, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

# Visa flow small at bottom center
if os.path.exists(FLOW_VISA):
    slide.shapes.add_picture(FLOW_VISA, Inches(4.0), Inches(5.9), Inches(3.5), Inches(1.3))

add_rounded_rect(slide, Inches(4.3), Inches(5.55), Inches(2.8), Inches(0.35), ACCENT_BLUE)
add_text(slide, Inches(4.4), Inches(5.57), Inches(2.6), Inches(0.3),
         "Visa Assistant Mapping", font_size=11, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_footer(slide, 6, total=11)

# ============================================================
# SLIDE 5: SCREEN LIST (Key Screens)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.6),
         "3. Screen List - Key UI Elements", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(0.9), Inches(3), Inches(0.04), ACCENT_BLUE)

# Screen categories with counts
screens_data = [
    ("Home / Dashboard", [
        "Flight Delay Banner - Real-time status",
        "Scan Passport - Camera trigger",
        "Check-in Button - Mobile check-in",
        "Show Boarding Pass - Digital QR",
        "Miles Earned Card - 42,850 miles tracked",
        "5-Tab Bottom Navigation Bar"
    ]),
    ("Visa Module", [
        "Passport Status Dashboard",
        "Destination Search Bar",
        "Visa-Free / e-Visa / Required Categories",
        "Start Application Button",
        "Documents Checklist",
        "Consular Location Map Widget"
    ]),
    ("Trips & Bookings", [
        "Journey Timeline - Chronological view",
        "Flight/Hotel/Activity Cards",
        "QR Ticket Access",
        "Add to Trip Selection Hub",
        "Accommodation Management"
    ]),
    ("Community & Analytics", [
        "Trending Destinations Feed",
        "Travel Stories & Blog Posts",
        "Buddy Search & Connect",
        "Spending Analysis Charts",
        "Global Footprint Progress",
        "Budget Management Tools"
    ]),
]

x = Inches(0.5)
for title, items in screens_data:
    # Title
    hdr = add_rounded_rect(slide, x, Inches(1.3), Inches(3), Inches(0.5), ACCENT_BLUE)
    tf = hdr.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Items
    y = Inches(1.95)
    for item in items:
        row = add_rounded_rect(slide, x, y, Inches(3), Inches(0.38),
                                RGBColor(0x14, 0x3A, 0x5C))
        add_text(slide, x + Inches(0.15), y + Inches(0.04), Inches(2.7), Inches(0.3),
                 item, font_size=10, color=MED_GRAY)
        y += Inches(0.42)

    x += Inches(3.15)

# Bottom highlight
add_rounded_rect(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.55),
                  RGBColor(0x14, 0x3A, 0x5C))
add_text(slide, Inches(1.3), Inches(5.87), Inches(10.7), Inches(0.4),
         "12 documented screens  |  100+ UI elements mapped  |  Every button, card, and navigation element catalogued",
         font_size=13, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Stats at bottom
stats = [("12", "Screens\nDocumented"), ("100+", "UI Elements\nMapped"),
         ("5", "Navigation\nModules"), ("3", "Task Flows\nAnalyzed")]
x = Inches(2)
for val, label in stats:
    add_text(slide, x, Inches(6.5), Inches(2), Inches(0.3),
             val, font_size=28, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(6.9), Inches(2), Inches(0.4),
             label, font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    x += Inches(2.5)

add_footer(slide, 7, total=11)

# ============================================================
# SLIDE 8: SCREEN DESIGNS - Dashboard & Flights
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.6),
         "4. Screen Designs - Dashboard & Flights", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(0.9), Inches(3), Inches(0.04), ACCENT_BLUE)

# Screenshots: Dashboard (SS[0]) and Check-in (SS[1])
if os.path.exists(SS[0]):
    slide.shapes.add_picture(SS[0], Inches(0.4), Inches(1.2), Inches(6.2), Inches(5.8))
if os.path.exists(SS[1]):
    slide.shapes.add_picture(SS[1], Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.8))

add_footer(slide, 8, total=11)

# ============================================================
# SLIDE 9: SCREEN DESIGNS - Visa & Trips
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.6),
         "4. Screen Designs - Visa Assistant & Trips", font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(0.9), Inches(3), Inches(0.04), ACCENT_BLUE)

# Visa (SS[2]) and Trips (SS[3])
if os.path.exists(SS[2]):
    slide.shapes.add_picture(SS[2], Inches(0.4), Inches(1.2), Inches(6.2), Inches(5.8))
if os.path.exists(SS[3]):
    slide.shapes.add_picture(SS[3], Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.8))

add_footer(slide, 9, total=11)

# ============================================================
# SLIDE 10: SCREEN DESIGNS - Planning, Community & Analytics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.3), Inches(10), Inches(0.6),
         "4. Screen Designs - Planning, Community & Analytics",
         font_size=32, color=WHITE, bold=True)
add_shape(slide, Inches(1), Inches(0.9), Inches(3), Inches(0.04), ACCENT_BLUE)

# Trip Planning (SS[4]), Community (SS[5]), Analytics (SS[6])
if os.path.exists(SS[4]):
    slide.shapes.add_picture(SS[4], Inches(0.2), Inches(1.2), Inches(4.2), Inches(5.8))
if os.path.exists(SS[5]):
    slide.shapes.add_picture(SS[5], Inches(4.6), Inches(1.2), Inches(4.2), Inches(5.8))
if os.path.exists(SS[6]):
    slide.shapes.add_picture(SS[6], Inches(8.9), Inches(1.2), Inches(4.2), Inches(5.8))

add_footer(slide, 10, total=11)

# ============================================================
# SLIDE 11: THANK YOU / Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

# Big thank you
add_text(slide, Inches(0), Inches(2.0), W, Inches(1),
         "Thank You", font_size=52, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(3.2), W, Inches(0.6),
         "Questions & Discussion", font_size=24, color=ACCENT_BLUE,
         alignment=PP_ALIGN.CENTER)

# Divider
add_shape(slide, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.04), ACCENT_BLUE)

# Team
add_text(slide, Inches(0), Inches(4.5), W, Inches(0.4),
         "Ozde Avdan  |  Nilsu Calliogullari  |  Deniz Tas  |  Fatih Bilal Yilmaz",
         font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(5.0), W, Inches(0.4),
         "IE48L  -  Human Computer Interaction Design  -  30.03.2026",
         font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# Logo
if os.path.exists(LOGO):
    slide.shapes.add_picture(LOGO, Inches(6.0), Inches(5.5), Inches(1.2), Inches(1.2))

add_footer(slide, 11, total=11)

# ============================================================
# SAVE
# ============================================================
prs.save(OUT)
print(f"PPTX saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
